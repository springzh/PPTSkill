#!/usr/bin/env python3
"""Build a new .pptx from a template by selecting slides and replacing text.

The script preserves all layout, fonts, colors, images and animations of the
source template. It only:
  1. Keeps the slides listed in ``selected_slides`` (in that order),
  2. Replaces the text of the runs identified by ``edits``.

USAGE
    python3 build_pptx.py <template.pptx> <edits.json> <output.pptx>
                         [--detail detail.json]   # for slot_id resolution
                         [--strict]               # fail if expected_text mismatch
                         [--dry-run]              # show planned edits, do nothing

EDITS JSON SCHEMA
{
  "template_slug": "simple-business",                 # optional, informational
  "selected_slides": [1, 2, 4, 7, 14, 16],            # 1-indexed in template
  "edits": [
    {
      "slide": 1,                                      # 1-indexed in template
      "slot_id": "title_en",                           # resolved via detail.json
      "new_text": "Annual Review"
    },
    {
      "slide": 1,
      "address": {"shape_id": 46, "paragraph": 1, "run": 0},  # explicit
      "expected_text": "工作计划模板",
      "new_text": "2026 年度工作计划"
    }
  ]
}

ADDRESS FORMS
    {"shape_id": <int>, "paragraph": <int>, "run": <int>}   # exact run
    {"shape_id": <int>, "paragraph": <int>}                  # whole paragraph
                                                              # (run 0 keeps fmt,
                                                              # others cleared)
"""
from __future__ import annotations
import argparse
import copy
import json
import math
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

SLIDE_R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


# ------------------------------------------------------------------
# Slot resolution via detail.json
# ------------------------------------------------------------------
def load_slot_index(detail_path: Path | None) -> dict:
    """Return {(slide_number, slot_id): slot_dict}.

    slot_dict includes address + capacity metadata (max_chars, chars_per_line,
    max_lines, wrap, autofit, font_size_pt, level, role) when present so the
    overflow lint can use them.
    """
    if detail_path is None or not detail_path.exists():
        return {}
    detail = json.loads(detail_path.read_text(encoding="utf-8"))
    index: dict = {}
    for page in detail.get("pages", []):
        slide_number = page["slide_number"]
        for slot in page.get("text_slots", []):
            slot.setdefault("expected_text", slot.get("current_text"))
            index[(slide_number, slot["slot_id"])] = slot
    return index


# ------------------------------------------------------------------
# Registry loading for multi-source mode
# ------------------------------------------------------------------
def load_registry(registry_path: Path) -> dict[str, dict]:
    """Load slide-registry.json. Returns {slide_id: slide_entry}."""
    if not registry_path.exists():
        raise FileNotFoundError(f"Registry not found: {registry_path}")
    data = json.loads(registry_path.read_text(encoding="utf-8"))
    index: dict[str, dict] = {}
    for slide in data.get("slides", []):
        index[slide["slide_id"]] = slide
    return index


def load_style(styles_dir: Path, style_slug: str) -> dict:
    """Load a style definition from styles/<slug>.json."""
    style_path = styles_dir / f"{style_slug}.json"
    if not style_path.exists():
        raise FileNotFoundError(f"Style not found: {style_path}")
    return json.loads(style_path.read_text(encoding="utf-8"))


def load_all_styles(styles_dir: Path) -> dict[str, dict]:
    """Load all style definitions. Returns {slug: style_dict}."""
    styles: dict[str, dict] = {}
    if not styles_dir.is_dir():
        return styles
    for style_path in sorted(styles_dir.glob("*.json")):
        style = json.loads(style_path.read_text(encoding="utf-8"))
        styles[style["slug"]] = style
    return styles


def resolve_slides(edits_spec: dict, registry: dict[str, dict]) -> tuple[list[dict], set[str]]:
    """Resolve slide_id entries to physical locations.

    Returns: (resolved_slides, set_of_source_pptx_paths)
    """
    resolved: list[dict] = []
    sources: set[str] = set()
    for entry in edits_spec.get("slides", []):
        sid = entry["slide_id"]
        if sid not in registry:
            raise KeyError(f"slide_id {sid!r} not found in registry")
        slide = registry[sid]
        resolved.append({
            "slide_id": sid,
            "source_pptx": slide["source_pptx"],
            "source_slide": slide["source_slide"],
            "style_family": slide["style_family"],
            "detail_ref": slide.get("detail_ref"),
        })
        sources.add(slide["source_pptx"])
    return resolved, sources


def build_merged_slot_index(resolved_slides: list[dict]) -> dict:
    """Build a merged slot index from multiple detail.json files.

    Returns: {(slide_id, slot_id): slot_metadata}
    """
    index: dict = {}
    detail_cache: dict[str, dict] = {}
    for entry in resolved_slides:
        ref = entry.get("detail_ref")
        if not ref:
            continue
        if ref not in detail_cache:
            detail_cache[ref] = json.loads(Path(ref).read_text(encoding="utf-8"))
        detail = detail_cache[ref]
        for page in detail.get("pages", []):
            if page["slide_number"] != entry["source_slide"]:
                continue
            for slot in page.get("text_slots", []):
                slot.setdefault("expected_text", slot.get("current_text"))
                index[(entry["slide_id"], slot["slot_id"])] = slot
            break
    return index


# ------------------------------------------------------------------
# Overflow lint (text-fits-the-box check)
# ------------------------------------------------------------------
def _visual_width(s: str) -> float:
    w = 0.0
    for c in s:
        if "\u4e00" <= c <= "\u9fff" or "\u3000" <= c <= "\u303f" or "\uff00" <= c <= "\uffef":
            w += 1.0
        elif c == " ":
            w += 0.35
        elif c.isascii():
            w += 0.5
        else:
            w += 0.8
    return w


def check_overflow(new_text: str, meta: dict) -> tuple[bool, str]:
    """Return (fits, message). Uses chars_per_line / max_lines from detail.json.

    'fits' is True when the text is expected to stay within the box. Boxes with
    autofit (PowerPoint shrink-to-fit) are always treated as fitting (soft).
    """
    if meta.get("capacity_unknown"):
        return True, ""  # geometry unreliable for this box
    cpl = meta.get("chars_per_line")
    max_lines = meta.get("max_lines")
    cap = meta.get("max_chars")   # already includes the calibration tolerance
    if not cpl or not max_lines or not cap:
        return True, ""  # no capacity data → can't lint
    total_vw = _visual_width(new_text.replace("\n", ""))
    # primary test: visual width vs the (tolerance-adjusted) capacity budget
    if total_vw <= cap:
        return True, ""
    # estimate needed lines for a helpful message
    needed = sum(max(1, math.ceil(_visual_width(seg) / cpl))
                 for seg in (new_text.split("\n") or [new_text]))
    sz = meta.get("font_size_pt")
    msg = (f"视觉宽度 {total_vw:.0f} > 容量 {cap} "
           f"(约需 {needed} 行 / 可容 {max_lines} 行, {sz}pt, 每行≈{cpl}字)")
    if meta.get("autofit"):
        return True, "AUTOFIT " + msg + " — PowerPoint 会自动缩字"
    return False, msg


# ------------------------------------------------------------------
# Shape lookup (recurses into groups)
# ------------------------------------------------------------------
def find_shape(shapes, shape_id: int):
    for shape in shapes:
        if shape.shape_id == shape_id:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found = find_shape(shape.shapes, shape_id)
            if found is not None:
                return found
    return None


# ------------------------------------------------------------------
# Text mutation that preserves run-level formatting
# ------------------------------------------------------------------
def apply_edit_to_shape(shape, address: dict, new_text: str, expected: str | None,
                       strict: bool) -> tuple[bool, str]:
    if not shape.has_text_frame:
        return False, f"shape {shape.shape_id} has no text frame"
    tf = shape.text_frame
    pi = address.get("paragraph", 0)
    if pi >= len(tf.paragraphs):
        return False, f"paragraph {pi} out of range (have {len(tf.paragraphs)})"
    para = tf.paragraphs[pi]
    runs = list(para.runs)
    if not runs:
        return False, f"paragraph {pi} has no runs"

    ri = address.get("run")
    if ri is None:
        current = "".join(r.text for r in runs)
        if expected is not None and current != expected and strict:
            return False, (
                f"expected_text mismatch at paragraph level: "
                f"have {current!r}, expected {expected!r}"
            )
        runs[0].text = new_text
        for r in runs[1:]:
            r.text = ""
        return True, f"replaced paragraph (kept run 0 formatting): {current!r} -> {new_text!r}"

    if ri >= len(runs):
        return False, f"run {ri} out of range (have {len(runs)})"
    current = runs[ri].text
    if expected is not None and current != expected and strict:
        return False, (
            f"expected_text mismatch: have {current!r}, expected {expected!r}"
        )
    runs[ri].text = new_text
    return True, f"replaced run: {current!r} -> {new_text!r}"


# ------------------------------------------------------------------
# Slide pruning: keep only selected slides, in order
# ------------------------------------------------------------------
def prune_slides(prs, selected_1indexed: list[int]) -> None:
    sldIdLst = prs.slides._sldIdLst
    sld_ids = list(sldIdLst)
    total = len(sld_ids)
    keep_idx0 = [i - 1 for i in selected_1indexed]
    for i in keep_idx0:
        if i < 0 or i >= total:
            raise ValueError(f"selected slide {i + 1} out of range (1..{total})")

    # Build new order list of element references.
    new_order = [sld_ids[i] for i in keep_idx0]

    # Remove all entries, then re-append in the new order. Keeping the same
    # element instances preserves their r:id and avoids orphaning relationships.
    for sld in list(sldIdLst):
        sldIdLst.remove(sld)
    for sld in new_order:
        sldIdLst.append(sld)


# ------------------------------------------------------------------
# Slide cloning for multi-source assembly
# ------------------------------------------------------------------
def _rewrite_rids(element, rid_map: dict[str, str]) -> None:
    """Recursively rewrite rId attributes in copied XML to match new relationships."""
    for attr_name in (
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed",
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link",
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}picture",
    ):
        for el in element.iter():
            val = el.get(attr_name)
            if val and val in rid_map:
                el.set(attr_name, rid_map[val])


def _clone_slide(src_slide, dst_prs):
    """Clone a slide from a source presentation into a destination presentation.

    Copies the slide XML, its relationships, and returns the new slide.
    """
    blank_layout = dst_prs.slide_layouts[6]  # blank layout
    dst_slide = dst_prs.slides.add_slide(blank_layout)

    src_elem = src_slide._element
    dst_elem = dst_slide._element

    # Remove all existing shapes from the blank slide
    NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    sp_tree = dst_elem.find(f".//{{{NS}}}spTree")
    if sp_tree is not None:
        for child in list(sp_tree):
            sp_tree.remove(child)
    else:
        # spTree not found as descendant; try as direct child or create one
        sp_tree = dst_elem.find(f"{{{NS}}}cSld/{{{NS}}}spTree")
        if sp_tree is None:
            cSld = dst_elem.find(f"{{{NS}}}cSld")
            if cSld is None:
                cSld = etree.SubElement(dst_elem, f"{{{NS}}}cSld")
            sp_tree = etree.SubElement(cSld, f"{{{NS}}}spTree")
        else:
            for child in list(sp_tree):
                sp_tree.remove(child)

    # Copy shapes from source
    src_sp_tree = src_elem.find(f".//{{{NS}}}spTree")
    if src_sp_tree is not None:
        for child in list(src_sp_tree):
            imported = copy.deepcopy(child)
            sp_tree.append(imported)

    # Clone relationships
    src_part = src_slide.part
    dst_part = dst_slide.part
    rid_map: dict[str, str] = {}

    for rel in src_part.rels.values():
        if rel.is_external:
            new_rid = dst_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            try:
                _blob = rel.target_part.blob
            except Exception:
                continue
            new_rid = dst_part.rels.get_or_add(rel.reltype, rel.target_part)
        rid_map[rel.rId] = new_rid

    if rid_map:
        _rewrite_rids(dst_elem, rid_map)

    return dst_slide


def assemble_slides(resolved_slides: list[dict], dst_prs) -> list[tuple[int, str]]:
    """Copy slides from multiple source .pptx files into destination presentation.

    Returns: list of (slide_index, style_family) for style application phase.
    """
    source_cache: dict[str, Presentation] = {}
    slide_meta: list[tuple[int, str]] = []

    for entry in resolved_slides:
        src_path = entry["source_pptx"]
        if src_path not in source_cache:
            source_cache[src_path] = Presentation(src_path)
        src_prs = source_cache[src_path]
        src_slide = src_prs.slides[entry["source_slide"] - 1]

        _clone_slide(src_slide, dst_prs)
        slide_meta.append((len(dst_prs.slides) - 1, entry["style_family"]))

    return slide_meta


# ------------------------------------------------------------------
# Style application (color + font remapping)
# ------------------------------------------------------------------
def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert '#1E3A5F' → RGBColor(0x1E, 0x3A, 0x5F)."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _rgb_to_hex(rgb) -> str:
    """Convert RGBColor → '#1E3A5F'."""
    return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def _try_rgb(value):
    """Safely extract hex string from a color value, or return None."""
    try:
        if value is None:
            return None
        return _rgb_to_hex(value)
    except Exception:
        return None


def build_color_map(source_style: dict, target_style: dict) -> dict[str, str]:
    """Build a hex→hex remap table from source style to target style."""
    src_cmap = source_style.get("color_map", {})
    tgt_palette = target_style.get("palette", {})
    result: dict[str, str] = {}
    for hex_val, role in src_cmap.items():
        hex_upper = hex_val.upper()
        if role in tgt_palette:
            result[hex_upper] = tgt_palette[role].upper()
        else:
            for fallback in ("accent", "secondary", "primary"):
                if fallback in tgt_palette and fallback != role:
                    result[hex_upper] = tgt_palette[fallback].upper()
                    break
    return result


def _collect_all_shapes(slide):
    """Recursively yield all shapes including those inside groups."""
    for shape in slide.shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _collect_all_shapes(shape)


def remap_shape_colors(shape, color_map: dict[str, str]) -> None:
    """Remap fill, line, and font colors on a shape if they match the color_map."""
    # Text runs
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color and run.font.color.rgb:
                    current = _try_rgb(run.font.color.rgb)
                    if current and current in color_map:
                        run.font.color.rgb = _hex_to_rgb(color_map[current])

    # Solid fill
    try:
        fill = shape.fill
        if fill.type is not None:
            if hasattr(fill, "fore_color") and fill.fore_color.rgb:
                current = _try_rgb(fill.fore_color.rgb)
                if current and current in color_map:
                    fill.fore_color.rgb = _hex_to_rgb(color_map[current])
            # Gradient stops
            if hasattr(fill, "gradient_stops"):
                for stop in fill.gradient_stops:
                    if stop.color.rgb:
                        current = _try_rgb(stop.color.rgb)
                        if current and current in color_map:
                            stop.color.rgb = _hex_to_rgb(color_map[current])
    except Exception:
        pass

    # Line color
    try:
        if hasattr(shape, "line") and shape.line.fill.type is not None:
            ln = shape.line
            if hasattr(ln, "color") and ln.color.rgb:
                current = _try_rgb(ln.color.rgb)
                if current and current in color_map:
                    ln.color.rgb = _hex_to_rgb(color_map[current])
    except Exception:
        pass


def remap_shape_fonts(shape, source_font_map: dict[str, str],
                      target_fonts: dict[str, str]) -> None:
    """Remap font names on a shape according to font_map → target_fonts."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            current_font = run.font.name
            if not current_font:
                continue
            role = source_font_map.get(current_font)
            if role:
                new_font = target_fonts.get(role)
                if new_font and new_font != current_font:
                    run.font.name = new_font


def apply_style(prs, target_style_slug: str, slide_meta: list[tuple[int, str]],
                styles: dict[str, dict]) -> None:
    """Apply target style to assembled slides via color/font remapping."""
    target = styles.get(target_style_slug)
    if target is None:
        print(f"WARN: target style {target_style_slug!r} not found, skipping style",
              file=sys.stderr)
        return

    for slide_idx, source_family in slide_meta:
        source = styles.get(source_family)
        if source is None:
            continue
        if source["slug"] == target["slug"]:
            continue  # Same family — no remapping needed

        color_map = build_color_map(source, target)
        src_font_map = source.get("font_map", {})
        tgt_fonts = target.get("fonts", {})

        slide = prs.slides[slide_idx]
        for shape in _collect_all_shapes(slide):
            if color_map:
                remap_shape_colors(shape, color_map)
            if src_font_map and tgt_fonts:
                remap_shape_fonts(shape, src_font_map, tgt_fonts)


# ------------------------------------------------------------------
# Multi-source driver
# ------------------------------------------------------------------
def run_multi_source(edits_spec: dict, args) -> int:
    """Execute multi-source mode: resolve → assemble → edit → style."""
    if args.registry is None:
        print("ERROR: multi-source mode requires --registry slide-registry.json",
              file=sys.stderr)
        return 2

    registry = load_registry(args.registry)
    resolved, sources = resolve_slides(edits_spec, registry)
    target_style = edits_spec.get("style")

    # Phase 1: Assemble slides from source(s)
    prs = Presentation()
    first_source = list(sources)[0] if sources else None
    if first_source:
        src_prs = Presentation(first_source)
        prs.slide_width = src_prs.slide_width
        prs.slide_height = src_prs.slide_height

    slide_meta = assemble_slides(resolved, prs)

    # Phase 2: Resolve edits (slot_id via merged detail.json index)
    slot_index = build_merged_slot_index(resolved)
    raw_edits = list(edits_spec.get("edits") or [])

    planned: list[dict] = []
    for e in raw_edits:
        slide_id = e.get("slide_id")
        if slide_id is None:
            print(f"ERROR: edit {e} missing slide_id", file=sys.stderr)
            return 2
        output_idx = None
        for idx, entry in enumerate(resolved):
            if entry["slide_id"] == slide_id:
                output_idx = idx
                break
        if output_idx is None:
            print(f"ERROR: slide_id {slide_id!r} not in slides list", file=sys.stderr)
            return 2

        if "address" in e:
            addr = dict(e["address"])
            expected = e.get("expected_text") or addr.get("expected_text")
            meta = None
        else:
            slot_id = e.get("slot_id")
            if slot_id is None:
                print(f"ERROR: edit {e} missing both address and slot_id",
                      file=sys.stderr)
                return 2
            key = (slide_id, slot_id)
            if key not in slot_index:
                print(f"ERROR: slot_id {slot_id!r} not found for slide {slide_id}",
                      file=sys.stderr)
                return 2
            meta = slot_index[key]
            addr = dict(meta["address"])
            expected = e.get("expected_text") or meta.get("expected_text")
        planned.append({
            "output_slide": output_idx,
            "address": addr,
            "expected": expected,
            "new_text": e["new_text"],
            "meta": meta,
            "raw": e,
        })

    if args.dry_run:
        for p in planned:
            print(f"slide {p['output_slide']:>3} addr={p['address']} "
                  f"-> {p['new_text']!r}")
        return 0

    # Phase 3: Apply text edits
    overflow_issues: list[str] = []
    for p in planned:
        slide = prs.slides[p["output_slide"]]
        shape_id = p["address"].get("shape_id")
        shape = find_shape(slide.shapes, shape_id)
        if shape is None:
            msg = f"slide {p['output_slide']}: shape_id {shape_id} not found"
            if args.strict:
                print(f"ERROR: {msg}", file=sys.stderr)
                return 3
            print(f"WARN:  {msg}", file=sys.stderr)
            continue
        ok, info = apply_edit_to_shape(
            shape, p["address"], p["new_text"], p["expected"], args.strict
        )
        prefix = "OK  " if ok else ("ERR " if args.strict else "WARN")
        print(f"{prefix} output slide {p['output_slide']:>3} shape {shape_id}: {info}")
        if not ok and args.strict:
            return 4

        if not args.no_lint and p["meta"]:
            fits, omsg = check_overflow(p["new_text"], p["meta"])
            if omsg and not fits:
                sid = p["meta"].get("slot_id", "?")
                line = (f"slide {p['output_slide']:>3} {sid}: 文字过长 {omsg} "
                        f"-> {p['new_text'][:24]!r}")
                overflow_issues.append(line)
                print(f"OVERFLOW {line}", file=sys.stderr)

    if overflow_issues:
        print(f"\n{len(overflow_issues)} 处文字可能出框:", file=sys.stderr)
        for line in overflow_issues:
            print(f"  - {line}", file=sys.stderr)
        if args.strict:
            print("\n--strict: 因出框风险拒绝保存。", file=sys.stderr)
            return 5

    # Phase 4: Apply style
    if target_style and args.styles:
        all_styles = load_all_styles(args.styles)
        apply_style(prs, target_style, slide_meta, all_styles)

    prs.save(args.output)
    note = (f"\nSaved {args.output} with {len(resolved)} slides "
            f"from {len(sources)} source(s) and {len(planned)} edits")
    if overflow_issues:
        note += f"  (⚠️ {len(overflow_issues)} 处可能出框，见上)"
    print(note)
    return 0


def run(args) -> int:
    edits_spec = json.loads(args.edits.read_text(encoding="utf-8"))
    selected = list(edits_spec.get("selected_slides") or [])
    raw_edits = list(edits_spec.get("edits") or [])
    slot_index = load_slot_index(args.detail)

    prs = Presentation(args.template)
    n_slides = len(prs.slides)
    if not selected:
        selected = list(range(1, n_slides + 1))
    for s in selected:
        if s < 1 or s > n_slides:
            print(f"ERROR: selected slide {s} out of range 1..{n_slides}", file=sys.stderr)
            return 2

    # Resolve and group edits by slide_number (in TEMPLATE order, pre-prune)
    planned: list[dict] = []
    for e in raw_edits:
        slide_num = e["slide"]
        meta = None
        if "address" in e:
            addr = dict(e["address"])
            expected = e.get("expected_text") or addr.get("expected_text")
        else:
            slot_id = e.get("slot_id")
            if slot_id is None:
                print(f"ERROR: edit {e} missing both address and slot_id", file=sys.stderr)
                return 2
            key = (slide_num, slot_id)
            if key not in slot_index:
                print(f"ERROR: slot_id {slot_id!r} not found for slide {slide_num} in detail.json",
                      file=sys.stderr)
                return 2
            meta = slot_index[key]
            addr = dict(meta["address"])
            expected = e.get("expected_text") or meta.get("expected_text")
        planned.append({
            "slide": slide_num,
            "address": addr,
            "expected": expected,
            "new_text": e["new_text"],
            "meta": meta,
            "raw": e,
        })

    if args.dry_run:
        for p in planned:
            print(f"slide {p['slide']:>3} addr={p['address']} -> {p['new_text']!r}")
        return 0

    # Apply edits BEFORE pruning so slide indices match the template
    overflow_issues: list[str] = []
    for p in planned:
        slide = prs.slides[p["slide"] - 1]
        shape_id = p["address"].get("shape_id")
        shape = find_shape(slide.shapes, shape_id)
        if shape is None:
            msg = f"slide {p['slide']}: shape_id {shape_id} not found"
            if args.strict:
                print(f"ERROR: {msg}", file=sys.stderr)
                return 3
            print(f"WARN:  {msg}", file=sys.stderr)
            continue
        ok, info = apply_edit_to_shape(
            shape, p["address"], p["new_text"], p["expected"], args.strict
        )
        prefix = "OK  " if ok else ("ERR " if args.strict else "WARN")
        print(f"{prefix} slide {p['slide']:>3} shape {shape_id}: {info}")
        if not ok and args.strict:
            return 4

        # Overflow lint (only when slot capacity metadata is available)
        if not args.no_lint and p["meta"]:
            fits, omsg = check_overflow(p["new_text"], p["meta"])
            if omsg and not fits:
                slot_id = p["meta"].get("slot_id", "?")
                line = (f"slide {p['slide']:>3} {slot_id}: 文字过长 {omsg} "
                        f"-> {p['new_text'][:24]!r}")
                overflow_issues.append(line)
                print(f"OVERFLOW {line}", file=sys.stderr)
            elif omsg and fits:  # autofit soft note
                print(f"note  slide {p['slide']:>3}: {omsg}", file=sys.stderr)

    if overflow_issues:
        print(f"\n{len(overflow_issues)} 处文字可能出框（请缩短文字到容量内，"
              f"保持字号不变以维持同级一致）：", file=sys.stderr)
        for line in overflow_issues:
            print(f"  - {line}", file=sys.stderr)
        if args.strict:
            print("\n--strict: 因出框风险拒绝保存。缩短上述文字后重试。", file=sys.stderr)
            return 5

    prune_slides(prs, selected)
    prs.save(args.output)
    note = f"\nSaved {args.output} with {len(selected)} slides and {len(planned)} edits"
    if overflow_issues:
        note += f"  (⚠️ {len(overflow_issues)} 处可能出框，见上)"
    print(note)
    return 0


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("arg1", type=Path,
                    help="template.pptx (legacy) or edits.json (multi-source)")
    ap.add_argument("arg2", type=Path, nargs="?", default=None,
                    help="edits.json (legacy) or output.pptx (multi-source)")
    ap.add_argument("arg3", type=Path, nargs="?", default=None,
                    help="output.pptx (legacy only)")
    ap.add_argument("--detail", type=Path, default=None)
    ap.add_argument("--registry", type=Path, default=None,
                    help="slide-registry.json (required for multi-source mode)")
    ap.add_argument("--styles", type=Path, default=None,
                    help="styles/ directory (for multi-source style application)")
    ap.add_argument("--no-lint", action="store_true",
                    help="disable the text-overflow check")
    ap.add_argument("--strict", action="store_true")
    ap.add_argument("--dry-run", action="store_true")
    args = ap.parse_args()

    # Auto-detect mode from the first positional argument
    # If arg1 is a .json with a "style" field → multi-source mode
    if args.arg1.suffix == ".json":
        spec = json.loads(args.arg1.read_text(encoding="utf-8"))
        if "style" in spec:
            if args.arg2 is None:
                print("ERROR: multi-source mode requires output path as second argument",
                      file=sys.stderr)
                return 2
            args.edits = args.arg1
            args.output = args.arg2
            return run_multi_source(spec, args)

    # Legacy mode: arg1=template.pptx, arg2=edits.json, arg3=output.pptx
    args.template = args.arg1
    if args.arg2 is None or args.arg3 is None:
        print("ERROR: legacy mode requires: template.pptx edits.json output.pptx",
              file=sys.stderr)
        return 2
    args.edits = args.arg2
    args.output = args.arg3

    if args.detail is None:
        guess = args.template.with_name("detail.json")
        if guess.exists():
            args.detail = guess
            print(f"(auto-detected detail.json: {guess})", file=sys.stderr)

    return run(args)


if __name__ == "__main__":
    sys.exit(main())
